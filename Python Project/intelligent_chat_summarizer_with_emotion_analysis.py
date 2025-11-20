"""
Intelligent Chat Summarizer with Emotion Analysis
Single-file Python project (CLI + optional Streamlit UI)

Contents:
- Preprocessing of chat transcripts
- Emotion detection per message using HuggingFace transformers pipeline
- Abstractive summarization using a T5 summarization pipeline with chunking
- Emotion analytics (distribution, timeline) and plots
- Export results: JSON, Markdown report, PPTX (slides)

How to run:
1) Create a virtualenv and install requirements:
   python -m venv venv
   source venv/bin/activate    # or venv\Scripts\activate on Windows
   pip install -r requirements.txt

2) Run the script on a sample chat file:
   python Intelligent_Chat_Summarizer_with_Emotion_Analysis.py --input chat.txt --output results

Files produced in the output folder:
- summary.txt        : generated abstractive summary
- emotions.csv       : per-message emotion labels and scores
- report.md          : markdown project report with visuals
- slides.pptx        : PPTX with summary + emotion charts
- emotion_timeline.png
- emotion_pie.png

Notes:
- This script uses HuggingFace transformer pipelines. Models are downloaded automatically when run if not present.
- For large chats, the summarizer chunks the chat to avoid model max-length limits, then summarizes chunks and summarizes again (hierarchical).

Recommended models (changeable in constants below):
- EMOTION_MODEL = "j-hartmann/emotion-english-distilroberta-base" (emotion classification)
- SUMMARIZATION_MODEL = "t5-small" (fine for small projects) or "t5-base" for better quality

"""

import os
import argparse
import json
import math
from pathlib import Path
from typing import List, Dict

import numpy as np
import pandas as pd
from transformers import pipeline, Pipeline
from sklearn.preprocessing import normalize

from pptx import Presentation
from pptx.util import Inches, Pt
import matplotlib.pyplot as plt


# ----------------------------
# Configuration
# ----------------------------
EMOTION_MODEL = "j-hartmann/emotion-english-distilroberta-base"
SUMMARIZATION_MODEL = "t5-small"
DEVICE = -1  # set to 0 if you have GPU (CUDA) and transformers configured

# Model pipeline batch sizes
EMOTION_BATCH = 32
SUMMARIZE_MAX_CHARS = 2500  # approx characters per chunk to feed into summarizer


# ----------------------------
# Helper functions
# ----------------------------

def read_chat_file(path: str) -> List[Dict[str, str]]:
    """Read a simple transcript file and return list of messages.

    Expected simple formats supported (auto-detected):
    1) WhatsApp exported txt (timestamp - Name: message) -> tries to parse
    2) Plain: one message per line

    Returns list of dicts: {"sender": str, "time": str, "text": str}
    """
    lines = Path(path).read_text(encoding='utf-8', errors='ignore').splitlines()
    messages = []
    for ln in lines:
        ln = ln.strip()
        if not ln:
            continue
        # Simple Whatsapp-ish detection: timestamp - Name: message
        if " - " in ln and ": " in ln.split(" - ", 1)[-1]:
            try:
                time_part, rest = ln.split(" - ", 1)
                sender, text = rest.split(": ", 1)
                messages.append({"sender": sender.strip(), "time": time_part.strip(), "text": text.strip()})
                continue
            except Exception:
                pass
        # fallback: treat whole line as text with unknown sender/time
        messages.append({"sender": "Unknown", "time": "Unknown", "text": ln})
    return messages


def chunk_texts(texts: List[str], max_chars: int) -> List[str]:
    """Join texts into chunks not exceeding max_chars (approx)."""
    chunks = []
    cur = []
    cur_len = 0
    for t in texts:
        if cur_len + len(t) + 1 > max_chars and cur:
            chunks.append(" ".join(cur))
            cur = [t]
            cur_len = len(t)
        else:
            cur.append(t)
            cur_len += len(t) + 1
    if cur:
        chunks.append(" ".join(cur))
    return chunks


# ----------------------------
# Pipelines
# ----------------------------

def get_emotion_pipeline(model_name: str, device: int = -1) -> Pipeline:
    return pipeline("text-classification", model=model_name, return_all_scores=True, device=device)


def get_summarizer_pipeline(model_name: str, device: int = -1) -> Pipeline:
    return pipeline("summarization", model=model_name, device=device)


# ----------------------------
# Core processing
# ----------------------------

def analyze_emotions(messages: List[Dict[str, str]], emotion_pipe: Pipeline) -> pd.DataFrame:
    texts = [m["text"] for m in messages]
    all_results = []
    for i in range(0, len(texts), EMOTION_BATCH):
        batch = texts[i:i + EMOTION_BATCH]
        results = emotion_pipe(batch)
        for j, res in enumerate(results):
            # res is list of dicts with 'label' and 'score'
            scores = {r['label']: r['score'] for r in res}
            all_results.append(scores)
    df = pd.DataFrame(all_results).fillna(0)
    df['text'] = texts
    df['sender'] = [m['sender'] for m in messages]
    df['time'] = [m['time'] for m in messages]
    # Top emotion
    df['dominant_emotion'] = df[[c for c in df.columns if c not in ('text','sender','time')]].idxmax(axis=1)
    return df


def summarize_chat(messages: List[Dict[str, str]], summarizer: Pipeline) -> str:
    texts = [m['text'] for m in messages]
    chunks = chunk_texts(texts, SUMMARIZE_MAX_CHARS)
    chunk_summaries = []
    for chunk in chunks:
        # summarizer may have max tokens limits; we pass chunk directly
        out = summarizer(chunk, max_length=120, min_length=30, do_sample=False)
        chunk_summaries.append(out[0]['summary_text'])
    # If multiple chunk summaries, summarize them again
    if len(chunk_summaries) == 1:
        return chunk_summaries[0]
    combined = " \n\n ".join(chunk_summaries)
    final = summarizer(combined, max_length=140, min_length=40, do_sample=False)
    return final[0]['summary_text']


# ----------------------------
# Reporting & visualization
# ----------------------------

def plot_emotion_pie(df: pd.DataFrame, outpath: str):
    counts = df['dominant_emotion'].value_counts()
    fig, ax = plt.subplots(figsize=(6,6))
    counts.plot.pie(autopct='%1.1f%%', ylabel='')
    ax.set_title('Emotion Distribution')
    fig.savefig(outpath)
    plt.close(fig)


def plot_emotion_timeline(df: pd.DataFrame, outpath: str):
    # Build timeline by message index
    labels = df.columns.drop(['text','sender','time','dominant_emotion'])
    timeline = df[list(labels)].to_numpy()
    fig, ax = plt.subplots(figsize=(10,4))
    ax.stackplot(range(len(df)), timeline.T)
    ax.set_xlim(0, len(df))
    ax.set_ylabel('Emotion score')
    ax.set_xlabel('Message index (time order)')
    ax.set_title('Emotion Timeline (stacked)')
    fig.savefig(outpath)
    plt.close(fig)


def save_pptx(summary: str, df: pd.DataFrame, output_path: str, pie_img: str, timeline_img: str):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Intelligent Chat Summarizer with Emotion Analysis"
    slide.placeholders[1].text = "Auto-generated report"

    # Summary slide
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Conversation Summary"
    body = slide.shapes.placeholders[1].text_frame
    body.text = summary

    # Emotion distribution
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Emotion Distribution"
    left = Inches(1)
    top = Inches(1.5)
    slide.shapes.add_picture(pie_img, left, top, width=Inches(8))

    # Emotion timeline
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Emotion Timeline"
    left = Inches(1)
    top = Inches(1.0)
    slide.shapes.add_picture(timeline_img, left, top, width=Inches(8))

    prs.save(output_path)


# ----------------------------
# Main CLI
# ----------------------------

def main():
    parser = argparse.ArgumentParser(description='Intelligent Chat Summarizer with Emotion Analysis')
    parser.add_argument('--input', required=True, help='Path to chat text file')
    parser.add_argument('--output', required=True, help='Output directory')
    parser.add_argument('--emotion_model', default=EMOTION_MODEL)
    parser.add_argument('--summarizer_model', default=SUMMARIZATION_MODEL)
    args = parser.parse_args()

    outdir = Path(args.output)
    outdir.mkdir(parents=True, exist_ok=True)

    print('Reading chat...')
    messages = read_chat_file(args.input)
    print(f'Loaded {len(messages)} messages')

    print('Loading emotion model... (this may download weights)')
    emotion_pipe = get_emotion_pipeline(args.emotion_model, device=DEVICE)

    print('Analyzing emotions...')
    df_em = analyze_emotions(messages, emotion_pipe)
    df_em.to_csv(outdir / 'emotions.csv', index=False)

    print('Plotting visuals...')
    pie_path = outdir / 'emotion_pie.png'
    timeline_path = outdir / 'emotion_timeline.png'
    plot_emotion_pie(df_em, str(pie_path))
    plot_emotion_timeline(df_em, str(timeline_path))

    print('Loading summarizer model...')
    summarizer = get_summarizer_pipeline(args.summarizer_model, device=DEVICE)

    print('Summarizing chat...')
    summary = summarize_chat(messages, summarizer)
    Path(outdir / 'summary.txt').write_text(summary, encoding='utf-8')

    print('Saving PPTX...')
    save_pptx(summary, df_em, str(outdir / 'slides.pptx'), str(pie_path), str(timeline_path))

    # Save report
    report_md = outdir / 'report.md'
    report_md.write_text(f"""
# Intelligent Chat Summarizer with Emotion Analysis

## Summary

{summary}

## Emotion Distribution

See emotion_pie.png and emotion_timeline.png

## Top emotional messages

""" + "\n\n".join(df_em.sort_values(by=df_em.columns[0], ascending=False).head(10)['text'].tolist()), encoding='utf-8')

    print('All done. Outputs written to', outdir)


if __name__ == '__main__':
    main()
